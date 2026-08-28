"""Generador de la presentacion de defensa del TFG Open CVN.

Construye docs/defensa/presentacion_tfg.pptx a partir del contenido de la
memoria (docs/memoria/chapters/*.tex) y de los diagramas ya generados en
docs/memoria/figs/. Ejecutar desde la raiz del repositorio:

    python docs/defensa/build_presentation.py

Requiere python-pptx y Pillow.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from PIL import Image
import os

ROOT = os.path.dirname(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))
FIGS = os.path.join(ROOT, "docs", "memoria", "figs")
OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "presentacion_tfg.pptx")

# ---------------------------------------------------------------------------
# Paleta y tipografia
# ---------------------------------------------------------------------------
PRIMARY = RGBColor(0x06, 0x5A, 0x82)      # azul profundo
SECONDARY = RGBColor(0x1C, 0x72, 0x93)    # verde azulado
DARK = RGBColor(0x11, 0x24, 0x3A)         # azul medianoche (fondos oscuros)
ACCENT = RGBColor(0xE0, 0x9A, 0x3E)       # ambar (acento puntual)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CARD_BG = RGBColor(0xF0, 0xF4, 0xF8)      # tarjeta clara con tinte azul
CARD_BG2 = RGBColor(0xE4, 0xEC, 0xF3)
TEXT_DARK = RGBColor(0x16, 0x22, 0x30)
TEXT_MUTED = RGBColor(0x5B, 0x6B, 0x7A)
TEXT_ON_DARK = RGBColor(0xE9, 0xF1, 0xF7)
TEXT_ON_DARK_MUTED = RGBColor(0xA9, 0xBF, 0xD1)
LINE_SOFT = RGBColor(0xD8, 0xE2, 0xEA)
AMBER_BG = RGBColor(0xFB, 0xF0, 0xDD)
TEAL_BG = RGBColor(0xE1, 0xEE, 0xF1)

HEAD_FONT = "Cambria"
BODY_FONT = "Calibri"

SLIDE_W_IN = 13.333
SLIDE_H_IN = 7.5

prs = Presentation()
prs.slide_width = Inches(SLIDE_W_IN)
prs.slide_height = Inches(SLIDE_H_IN)
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_slide(bg=WHITE):
    slide = prs.slides.add_slide(BLANK)
    bg_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = bg
    bg_shape.line.fill.background()
    bg_shape.shadow.inherit = False
    # send to back
    spTree = slide.shapes._spTree
    spTree.remove(bg_shape._element)
    spTree.insert(2, bg_shape._element)
    return slide


def add_rect(slide, x, y, w, h, fill=None, line_color=None, line_w=0.75, rounded=False, radius=0.08):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    if fill is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line_color is not None:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def add_oval(slide, x, y, d, fill=PRIMARY, line_color=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line_color is not None:
        shp.line.color.rgb = line_color
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def set_cell_text(shape, text, size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER, font=BODY_FONT):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.02)
    tf.margin_right = Inches(0.02)
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color


def add_text(slide, x, y, w, h, text, size=16, color=TEXT_DARK, bold=False, italic=False,
             font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=None,
             space_after=None, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if line_spacing:
            p.line_spacing = line_spacing
        if space_after:
            p.space_after = Pt(space_after)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.name = font
        r.font.color.rgb = color
    return tb


def _set_bullet(paragraph, color_hex, indent_in=0.28, char="●", scale=55):
    pPr = paragraph._p.get_or_add_pPr()
    pPr.set('marL', str(Emu(Inches(indent_in))))
    pPr.set('indent', str(-Emu(Inches(indent_in))))
    buClr = pPr.makeelement(qn('a:buClr'), {})
    srgb = buClr.makeelement(qn('a:srgbClr'), {'val': color_hex})
    buClr.append(srgb)
    buSzPct = pPr.makeelement(qn('a:buSzPct'), {'val': str(scale * 1000)})
    buFont = pPr.makeelement(qn('a:buFont'), {'typeface': 'Arial'})
    buChar = pPr.makeelement(qn('a:buChar'), {'char': char})
    pPr.append(buClr)
    pPr.append(buSzPct)
    pPr.append(buFont)
    pPr.append(buChar)


def add_bullets(slide, x, y, w, h, items, size=15, color=TEXT_DARK, font=BODY_FONT,
                 bullet_hex="065A82", space_after=10, anchor=MSO_ANCHOR.TOP, line_spacing=1.05):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        _set_bullet(p, bullet_hex)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            r1.text = lead
            r1.font.bold = True
            r1.font.size = Pt(size)
            r1.font.name = font
            r1.font.color.rgb = color
            r2 = p.add_run()
            r2.text = rest
            r2.font.size = Pt(size)
            r2.font.name = font
            r2.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(size)
            r.font.name = font
            r.font.color.rgb = color
    return tb


def add_image_fit(slide, path, x, y, max_w, max_h, align="center"):
    im = Image.open(path)
    aspect = im.size[0] / im.size[1]
    w = max_w
    h = w / aspect
    if h > max_h:
        h = max_h
        w = h * aspect
    if align == "center":
        left = x + (max_w - w) / 2
    else:
        left = x
    top = y + (max_h - h) / 2
    slide.shapes.add_picture(path, Inches(left), Inches(top), Inches(w), Inches(h))
    return left, top, w, h


def add_notes(slide, text):
    notes = slide.notes_slide
    notes.notes_text_frame.text = text


def kicker(slide, text, on_dark=False):
    color = ACCENT if on_dark else SECONDARY
    add_text(slide, 0.6, 0.42, 10.0, 0.35, text.upper(), size=13, bold=True,
             color=color, font=BODY_FONT)


def title(slide, text, y=0.72, size=30, on_dark=False, w=12.1):
    color = WHITE if on_dark else DARK
    add_text(slide, 0.6, y, w, 0.9, text, size=size, bold=True, color=color, font=HEAD_FONT,
              line_spacing=1.02)


def page_num(slide, n):
    add_text(slide, 12.55, 7.12, 0.6, 0.3, str(n), size=10, color=TEXT_MUTED, align=PP_ALIGN.RIGHT)


def circle_num(slide, x, y, d, number, fill=PRIMARY, text_color=WHITE, size=15):
    c = add_oval(slide, x, y, d, fill=fill)
    set_cell_text(c, str(number), size=size, color=text_color, bold=True)
    return c


def stat_card(slide, x, y, w, h, number, label, num_color=PRIMARY, bg=CARD_BG):
    add_rect(slide, x, y, w, h, fill=bg, rounded=True, radius=0.10)
    add_text(slide, x + 0.08, y + 0.10, w - 0.16, h - 0.62, number, size=30, bold=True,
              color=num_color, align=PP_ALIGN.CENTER, font=HEAD_FONT, anchor=MSO_ANCHOR.BOTTOM)
    add_text(slide, x + 0.08, y + h - 0.5, w - 0.16, 0.46, label, size=10.5, color=TEXT_MUTED,
              align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP, line_spacing=1.0)


def style_table(table, n_rows, n_cols, header_fill=PRIMARY, header_color=WHITE,
                 body_font_size=11.5, header_font_size=12, alt_fill=CARD_BG, body_color=TEXT_DARK,
                 align_first_left=True):
    table.first_row = False
    table.horz_banding = False
    for r in range(n_rows):
        for c in range(n_cols):
            cell = table.cell(r, c)
            cell.margin_left = Inches(0.09)
            cell.margin_right = Inches(0.09)
            cell.margin_top = Inches(0.045)
            cell.margin_bottom = Inches(0.045)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame
            tf.word_wrap = True
            for p in tf.paragraphs:
                p.alignment = PP_ALIGN.LEFT if (align_first_left and c == 0) else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = BODY_FONT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_fill
                for p in tf.paragraphs:
                    for run in p.runs:
                        run.font.bold = True
                        run.font.size = Pt(header_font_size)
                        run.font.color.rgb = header_color
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if r % 2 == 1 else alt_fill
                for p in tf.paragraphs:
                    for run in p.runs:
                        run.font.size = Pt(body_font_size)
                        run.font.color.rgb = body_color


def fill_table(table, rows):
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text_frame.text = str(val)


def add_flow_row(slide, x, y, total_w, h, items, fill=SECONDARY, text_color=WHITE, size=11.5,
                  arrow_w=0.34):
    n = len(items)
    arrow_total = arrow_w * (n - 1)
    avail = total_w - arrow_total
    lengths = [max(len(t), 10) for t in items]
    total_len = sum(lengths)
    widths = [avail * l / total_len for l in lengths]
    cx = x
    for i, (t, w) in enumerate(zip(items, widths)):
        add_rect(slide, cx, y, w, h, fill=fill, rounded=True, radius=0.5)
        add_text(slide, cx + 0.04, y, w - 0.08, h, t, size=size, bold=True, color=text_color,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_spacing=0.95)
        cx += w
        if i < n - 1:
            add_text(slide, cx, y, arrow_w, h, "→", size=16, bold=True, color=TEXT_MUTED,
                      align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
            cx += arrow_w
    return cx


# ---------------------------------------------------------------------------
# Slide 1 - Portada
# ---------------------------------------------------------------------------

def slide_01_portada():
    s = new_slide(bg=DARK)

    # motivo decorativo: capas apiladas (coherente con la arquitectura por capas)
    layer_colors = [PRIMARY, SECONDARY, ACCENT]
    bar_y = 6.78
    widths = [13.333, 9.4, 5.6]
    for i, (w, col) in enumerate(zip(widths, layer_colors)):
        add_rect(s, 0, bar_y + i * 0.24, w, 0.20, fill=col)

    add_text(s, 0.7, 0.62, 8, 0.4, "TRABAJO FIN DE GRADO · INTENSIFICACIÓN EN COMPUTACIÓN",
             size=12.5, bold=True, color=ACCENT, font=BODY_FONT)

    add_text(s, 0.7, 1.35, 11.6, 2.3,
             "Especificación de formato y desarrollo de herramientas para el "
             "procesamiento de currículum vítae en el ámbito académico universitario",
             size=30, bold=True, color=WHITE, font=HEAD_FONT, line_spacing=1.08)

    add_text(s, 0.7, 3.62, 11.4, 0.85,
             "Open CVN: una arquitectura computacional abierta para representar, "
             "validar, transformar y exportar currículos académicos",
             size=16.5, italic=True, color=TEXT_ON_DARK, font=BODY_FONT, line_spacing=1.05)

    add_text(s, 0.7, 4.75, 8, 0.4, "Carlos Martínez Jaén", size=17, bold=True, color=WHITE, font=HEAD_FONT)
    add_text(s, 0.7, 5.18, 8, 0.35, "Directores: Luis De la Ossa Jiménez · José Antonio Gámez Martín",
             size=12.5, color=TEXT_ON_DARK_MUTED)
    add_text(s, 0.7, 5.52, 8, 0.35, "Escuela Superior de Ingeniería Informática de Albacete · "
                                     "Universidad de Castilla-La Mancha",
             size=12.5, color=TEXT_ON_DARK_MUTED)
    add_text(s, 0.7, 5.86, 8, 0.35, "Julio de 2026", size=12.5, color=TEXT_ON_DARK_MUTED)

    add_notes(s,
              "Duracion objetivo: 30 s. Saludo al tribunal, presentacion personal breve y "
              "titulo del TFG. Frase puente: 'Voy a presentar Open CVN, una arquitectura "
              "computacional abierta para representar, validar, transformar y exportar "
              "curriculos academicos tomando CVN como referencia.' Recordar en voz baja el "
              "guion propio: objetivos, metodologia, contenido y conclusiones, en 25 minutos.")


# ---------------------------------------------------------------------------
# Slide 2 - Indice
# ---------------------------------------------------------------------------

def slide_02_indice():
    s = new_slide()
    kicker(s, "Guion de la defensa")
    title(s, "Índice")

    items = [
        "Motivación y contexto",
        "Objetivos del trabajo",
        "El ecosistema CVN y su análisis",
        "Metodología y arquitectura",
        "Implementación del pipeline",
        "Formato Open CVN y herramienta",
        "Evaluación y resultados",
        "Conclusiones y trabajo futuro",
    ]
    col_x = [0.7, 6.9]
    col_w = 5.5
    row_h = 0.86
    for i, it in enumerate(items):
        col = i // 4
        row = i % 4
        x = col_x[col]
        y = 1.95 + row * row_h
        circle_num(s, x, y, 0.5, i + 1, fill=PRIMARY if i % 2 == 0 else SECONDARY)
        add_text(s, x + 0.68, y + 0.02, col_w - 0.7, 0.5, it, size=16, color=TEXT_DARK,
                  bold=True, font=BODY_FONT, anchor=MSO_ANCHOR.MIDDLE)

    add_rect(s, 0.7, 5.95, 11.9, 0.9, fill=CARD_BG, rounded=True, radius=0.12)
    add_text(s, 0.95, 6.10, 11.4, 0.65,
              "25 minutos de exposición: objetivos, metodología, contenido y conclusiones del "
              "trabajo. Después, turno de preguntas del tribunal.",
              size=13, italic=True, color=TEXT_MUTED, anchor=MSO_ANCHOR.MIDDLE)

    add_notes(s,
              "Duracion objetivo: 40 s. Recorrer el indice en voz alta rapidamente, sin "
              "detenerse en cada punto. Insistir en que la estructura sigue lo que exige la "
              "normativa: objetivos, metodologia, contenido y conclusiones.")


# ---------------------------------------------------------------------------
# Slide 3 - Contexto y motivacion
# ---------------------------------------------------------------------------

def slide_03_contexto():
    s = new_slide()
    kicker(s, "1 · Motivación y contexto")
    title(s, "El coste no está en el formato, está en usarlo")

    add_bullets(s, 0.7, 1.95, 6.9, 4.6, [
        ("Norma de referencia. ", "FECYT define CVN como el formato común de presentación de "
         "datos curriculares de investigadores en España, y su interoperabilidad con bases de "
         "datos institucionales."),
        ("Uso más allá de la presentación. ", "El currículo se emplea como entrada en "
         "evaluaciones, convocatorias, solicitudes de ayudas, acreditaciones y sistemas de "
         "gestión de la investigación."),
        ("El formato existe, el problema persiste. ", "Disponer de una norma común no elimina "
         "el coste de elaborar, mantener y adaptar un currículo a distintos contextos."),
    ], size=15.5, space_after=18, line_spacing=1.08)

    add_rect(s, 7.95, 1.95, 4.65, 4.35, fill=AMBER_BG, rounded=True, radius=0.10)
    add_text(s, 8.25, 2.20, 4.05, 0.4, "EL PROBLEMA DE PARTIDA", size=12, bold=True, color=ACCENT)
    add_text(s, 8.25, 2.68, 4.05, 2.3,
              "“Pequeñas variaciones en la estructura o en el formato de salida obligan a "
              "tareas manuales repetitivas cada vez que un currículo se reutiliza en un "
              "contexto distinto.”",
              size=16, italic=True, color=TEXT_DARK, font=HEAD_FONT, line_spacing=1.18)
    add_text(s, 8.25, 5.55, 4.05, 0.6,
              "El uso automatizado de CVN sigue siendo costoso.", size=12.5, bold=True,
              color=PRIMARY)

    add_notes(s,
              "Duracion objetivo: 90 s. Explicar que CVN (FECYT) normaliza la presentacion "
              "curricular en Espana, pero que la automatizacion real sigue siendo costosa: "
              "actualizar, adaptar y reutilizar un curriculo exige trabajo manual. Esta es la "
              "motivacion de partida del TFG, antes de entrar en como funciona CVN por dentro.")


# ---------------------------------------------------------------------------
# Slide 4 - Como funciona CVN
# ---------------------------------------------------------------------------

def slide_04_como_funciona():
    s = new_slide()
    kicker(s, "1 · Motivación y contexto")
    title(s, "CVN no es un documento: es un ecosistema de artefactos")

    layers = [
        ("Manual de especificaciones (PDF)", "Significado funcional de cada campo, para lectura humana", PRIMARY),
        ("Manual estructurado (XML)", "El mismo manual, en forma parseable por software", SECONDARY),
        ("Modelo en árbol (XML)", "Enlaza cada código funcional con el XML real de un currículo", SECONDARY),
        ("Esquemas XSD", "Forma legal del XML y vocabulario de tablas auxiliares", DARK),
    ]
    x0, w0 = 0.7, 6.55
    y = 1.95
    rh, gap = 0.86, 0.14
    for label, desc, col in layers:
        add_rect(s, x0, y, w0, rh, fill=col, rounded=True, radius=0.10)
        add_text(s, x0 + 0.25, y + 0.10, w0 - 0.5, 0.36, label, size=14.5, bold=True,
                  color=WHITE, font=BODY_FONT)
        add_text(s, x0 + 0.25, y + 0.44, w0 - 0.5, 0.36, desc, size=11, color=TEXT_ON_DARK_MUTED,
                  font=BODY_FONT)
        y += rh + gap

    chip_y = y + 0.08
    add_text(s, x0, chip_y, w0, 0.3, "+ tres familias auxiliares del paquete oficial:",
              size=11.5, italic=True, color=TEXT_MUTED)
    chips = ["Catálogo de entidades", "Tablas de referencia / subtipos", "Tesauro multilingüe"]
    cx = x0
    for ch in chips:
        cw = 0.16 + 0.092 * len(ch)
        add_rect(s, cx, chip_y + 0.34, cw, 0.36, fill=CARD_BG2, rounded=True, radius=0.5)
        add_text(s, cx, chip_y + 0.34, cw, 0.36, ch, size=10.5, bold=True, color=PRIMARY,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        cx += cw + 0.14

    add_bullets(s, 7.65, 1.95, 4.95, 4.6, [
        ("Fuente única y oficial. ", "Todo el trabajo parte de la norma CVN y del paquete que "
         "FECYT distribuye para su implementación."),
        ("Qué significa vs. cómo se serializa. ", "El manual explica el significado de cada "
         "campo; el modelo en árbol y los XSD explican cómo se codifica en XML."),
        ("XML es el mecanismo, no el objeto. ", "XML sirve para representar e intercambiar los "
         "datos de CVN; no agota el contenido conceptual de la norma."),
    ], size=14, space_after=16, line_spacing=1.08)

    add_notes(s,
              "Duracion objetivo: 100 s. Sin entrar en estado del arte: solo explicar en que se "
              "basa CVN. Mostrar las cuatro capas del paquete oficial (manual PDF, manual XML, "
              "modelo en arbol, XSD) mas las tres familias auxiliares (entidades, tablas de "
              "referencia/subtipos, tesauro). Idea clave: CVN es un ecosistema de artefactos "
              "oficiales, XML es el mecanismo de representacion, no el unico objeto del trabajo.")


# ---------------------------------------------------------------------------
# Slide 5 - Objetivos
# ---------------------------------------------------------------------------

def slide_05_objetivos():
    s = new_slide()
    kicker(s, "2 · Objetivos")
    title(s, "Objetivo general")

    add_rect(s, 0.7, 1.85, 11.9, 1.55, fill=CARD_BG, rounded=True, radius=0.09)
    add_text(s, 0.98, 2.0, 0.5, 0.3, "OG", size=13, bold=True, color=ACCENT)
    add_text(s, 0.98, 2.28, 11.35, 1.02,
              "Diseñar e implementar una arquitectura computacional abierta para representar, "
              "validar, transformar, almacenar y exportar currículos académicos e "
              "investigadores tomando CVN como norma de referencia, sin que su representación "
              "XML ni la herramienta oficial condicionen por completo el modelo interno.",
              size=15, italic=True, color=TEXT_DARK, font=HEAD_FONT, line_spacing=1.14)

    groups = [
        ("OE1–OE2", "Análisis", "Analizar el ecosistema CVN y estudiar alternativas de "
         "representación y validación de datos curriculares."),
        ("OE3–OE6", "Generación y normalización", "Generar la capa estructural desde los XSD, "
         "normalizar metadatos, definir reglas semánticas y modelos de dominio."),
        ("OE7–OE9", "Formato y herramienta", "Definir el formato Open CVN JSON e implementar "
         "parseo, validación, almacenamiento, exportación e importación."),
        ("OE10", "Verificación", "Verificar el sistema mediante pruebas automatizadas y flujos "
         "reproducibles de extremo a extremo."),
    ]
    n = len(groups)
    gap = 0.16
    cw = (11.9 - gap * (n - 1)) / n
    x = 0.7
    y = 3.65
    ch = 2.75
    for code, name, desc in groups:
        add_rect(s, x, y, cw, ch, fill=WHITE, line_color=LINE_SOFT, line_w=1.1, rounded=True, radius=0.09)
        add_rect(s, x, y, cw, 0.62, fill=PRIMARY, rounded=True, radius=0.30)
        add_rect(s, x, y + 0.31, cw, 0.31, fill=PRIMARY)
        add_text(s, x, y, cw, 0.62, code, size=15, bold=True, color=WHITE,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD_FONT)
        add_text(s, x + 0.16, y + 0.78, cw - 0.32, 0.5, name, size=13.5, bold=True,
                  color=DARK, font=BODY_FONT, line_spacing=1.0)
        add_text(s, x + 0.16, y + 1.26, cw - 0.32, ch - 1.4, desc, size=10.8, color=TEXT_MUTED,
                  line_spacing=1.12)
        x += cw + gap

    add_text(s, 0.7, 6.62, 11.9, 0.5,
              "10 objetivos específicos, organizados en cuatro bloques alineados con los "
              "capítulos 3 a 7 de la memoria.",
              size=12, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_notes(s,
              "Duracion objetivo: 100 s. Leer el objetivo general con calma (esa frase es la "
              "columna vertebral de toda la defensa, volvera en conclusiones). Despues recorrer "
              "los 4 bloques de objetivos especificos sin detenerse en cada OE individual; "
              "mencionar que se retomaran uno a uno en las conclusiones con su grado de "
              "cumplimiento.")


# ---------------------------------------------------------------------------
# Slide 6 - Limitaciones del ecosistema CVN
# ---------------------------------------------------------------------------

def slide_06_limitaciones():
    s = new_slide()
    kicker(s, "3 · El ecosistema CVN")
    title(s, "El paquete oficial no es directamente un modelo interno")

    rows = [
        ["Limitación detectada", "Respuesta de diseño de Open CVN"],
        ["Discrepancias puntuales entre el XML canónico y su esquema de validación",
         "Tratar el XML como evidencia de origen y registrar la discrepancia de forma explícita"],
        ["Construcciones XSD que no se trasladan de forma perfecta a estructuras ejecutables "
         "(choice, cardinalidades, tipos genéricos)",
         "Confinar estas limitaciones a la capa estructural generada, sin exponerlas en el "
         "formato final"],
        ["Tablas auxiliares que no siempre son enumeraciones cerradas",
         "Evaluar la elegibilidad de enumeración cerrada a partir de evidencia concreta de "
         "cada tabla"],
        ["Referencias sin tabla equivalente y deriva de empaquetado entre familias auxiliares",
         "Conservar como no resueltas o parcialmente trazables, sin forzar una conversión con "
         "pérdida de información"],
    ]
    tbl_shape = s.shapes.add_table(len(rows), 2, Inches(0.7), Inches(1.95), Inches(11.9), Inches(4.55))
    table = tbl_shape.table
    table.columns[0].width = Inches(5.55)
    table.columns[1].width = Inches(6.35)
    fill_table(table, rows)
    style_table(table, len(rows), 2, body_font_size=12.5, header_font_size=13.5)

    add_text(s, 0.7, 6.68, 11.9, 0.45,
              "Ninguna limitación invalida la norma CVN: condicionan su uso directo como "
              "modelo interno de una herramienta.",
              size=11.5, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_notes(s,
              "Duracion objetivo: 90 s. Mostrar evidencia tecnica concreta sin dramatizar: son "
              "cuatro tipos de limitacion detectados al analizar el paquete oficial, cada una "
              "con una respuesta de diseno ya en Open CVN. Cerrar dejando claro que CVN sigue "
              "siendo valido para su proposito original de interoperabilidad documental.")


# ---------------------------------------------------------------------------
# Slide 7 - Propuesta Open CVN
# ---------------------------------------------------------------------------

def slide_07_propuesta():
    s = new_slide()
    kicker(s, "3 · El ecosistema CVN")
    title(s, "Open CVN: una capa trazable sobre el paquete oficial")

    add_image_fit(s, os.path.join(FIGS, "open_cvn_source_separation.png"),
                   0.7, 1.80, 11.9, 2.55)

    add_flow_row(s, 0.7, 4.55, 11.9, 0.66,
                 ["Paquete CVN", "Metadatos normalizados", "Política semántica",
                  "Modelo conceptual", "Open CVN JSON"],
                 fill=PRIMARY, size=12)

    add_text(s, 0.7, 5.55, 11.9, 0.9,
              "Cada capa añade una interpretación adicional sin descartar la relación con el "
              "paquete oficial: todo elemento del modelo final conserva evidencia verificable "
              "de la fuente CVN de la que procede.",
              size=14, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER, line_spacing=1.15)

    add_notes(s,
              "Duracion objetivo: 80 s. Presentar la figura de separacion entre CVN y Open CVN: "
              "el paquete oficial no se copia como modelo interno, se transforma de forma "
              "trazable en metadatos normalizados, luego politica semantica, luego modelo "
              "conceptual y finalmente el formato Open CVN JSON. Esta diapositiva es el puente "
              "hacia metodologia y arquitectura.")


# ---------------------------------------------------------------------------
# Slide 8 - Metodologia y herramientas
# ---------------------------------------------------------------------------

def slide_08_metodologia():
    s = new_slide()
    kicker(s, "4 · Metodología y arquitectura")
    title(s, "Cinco fases, con documentación continua de decisiones")

    phases = ["Análisis\ninicial", "Exploración de\nartefactos", "Diseño incremental\npor capas",
               "Implementación\nreproducible", "Verificación\nautomatizada"]
    n = len(phases)
    x0, x1 = 1.15, 12.2
    line_y = 2.32
    add_rect(s, x0, line_y - 0.012, x1 - x0, 0.03, fill=LINE_SOFT)
    step_w = (x1 - x0) / (n - 1)
    for i, ph in enumerate(phases):
        cx = x0 + step_w * i
        circle_num(s, cx - 0.28, line_y - 0.28, 0.56, i + 1,
                   fill=PRIMARY if i % 2 == 0 else SECONDARY)
        add_text(s, cx - 0.95, line_y + 0.42, 1.9, 0.65, ph, size=11.5, bold=True,
                  color=TEXT_DARK, align=PP_ALIGN.CENTER, line_spacing=1.0)

    add_text(s, 0.7, 3.55, 11.9, 0.4,
              "Herramientas: cada elección permite generar o validar artefactos de forma "
              "reproducible", size=13.5, bold=True, color=DARK, font=BODY_FONT)

    tools = [
        ("PY", "Python + uv"), ("XSD", "xsdata"), ("PYD", "Pydantic"), ("JS", "JSON Schema"),
        ("DB", "SQLite"), ("J2", "Jinja + LaTeX"), ("PT", "pytest + xdist"), ("GH", "Git + Actions"),
    ]
    n2 = len(tools)
    gap2 = 0.18
    cw2 = (11.9 - gap2 * (n2 - 1)) / n2
    x = 0.7
    y2 = 4.15
    for code, name in tools:
        add_oval(s, x + cw2 / 2 - 0.36, y2, 0.72, fill=CARD_BG2)
        add_text(s, x, y2, cw2, 0.72, code, size=13, bold=True, color=PRIMARY,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, font=HEAD_FONT)
        add_text(s, x - 0.1, y2 + 0.82, cw2 + 0.2, 0.7, name, size=10.5, color=TEXT_DARK,
                  bold=True, align=PP_ALIGN.CENTER, line_spacing=1.0)
        x += cw2 + gap2

    add_notes(s,
              "Duracion objetivo: 90 s. El proceso no fue un desarrollo de producto con "
              "sprints: cinco fases (analisis, exploracion de artefactos, diseno incremental, "
              "implementacion reproducible, verificacion) con documentacion continua de "
              "decisiones y limitaciones como criterio transversal. Repasar brevemente las "
              "herramientas clave: Python/uv, xsdata, Pydantic, JSON Schema, SQLite, Jinja, "
              "pytest, Git/GitHub Actions.")


# ---------------------------------------------------------------------------
# Slide 9 - Arquitectura general por capas
# ---------------------------------------------------------------------------

def slide_09_arquitectura():
    s = new_slide()
    kicker(s, "4 · Metodología y arquitectura")
    title(s, "Arquitectura general por capas")

    add_image_fit(s, os.path.join(FIGS, "open_cvn_layered_architecture.png"),
                   0.7, 1.85, 11.9, 2.05)

    blocks = [
        ("Fuente oficial CVN", "Frontera externa: única fuente formal del sistema"),
        ("Generación y normalización", "Bindings estructurales, metadatos y política semántica"),
        ("Dominio, conceptual y formato", "Modelos de dominio, modelo conceptual y JSON Schema"),
        ("Herramienta local", "Uso operativo: importación, versiones, exportación"),
    ]
    n = len(blocks)
    gap = 0.16
    cw = (11.9 - gap * (n - 1)) / n
    x = 0.7
    y = 4.35
    ch = 2.05
    colors = [DARK, PRIMARY, SECONDARY, ACCENT]
    for (name, desc), col in zip(blocks, colors):
        add_rect(s, x, y, cw, ch, fill=CARD_BG, rounded=True, radius=0.09)
        add_rect(s, x, y, cw, 0.10, fill=col, rounded=True, radius=0.5)
        add_text(s, x + 0.16, y + 0.28, cw - 0.32, 0.7, name, size=12.5, bold=True, color=DARK,
                  line_spacing=1.05)
        add_text(s, x + 0.16, y + 1.0, cw - 0.32, ch - 1.15, desc, size=10.5, color=TEXT_MUTED,
                  line_spacing=1.12)
        x += cw + gap

    add_notes(s,
              "Duracion objetivo: 90 s. Mostrar la figura de arquitectura por capas y explicar "
              "los cuatro bloques: fuente CVN (frontera externa), generacion y normalizacion "
              "(bindings + metadatos + politica semantica), dominio/modelo conceptual/formato, "
              "y herramienta local. Cada bloque corresponde a un modulo real del repositorio; "
              "el detalle de implementacion vendra en la siguiente diapositiva.")


# ---------------------------------------------------------------------------
# Slide 10 - Pipeline: generacion y normalizacion
# ---------------------------------------------------------------------------

def slide_10_pipeline():
    s = new_slide()
    kicker(s, "5 · Implementación del pipeline")
    title(s, "De los XSD oficiales a una política semántica trazable")

    steps = [
        ("1", "Generación estructural", "xsdata genera bindings Pydantic fieles al XML de CVN, "
         "sin interpretación semántica todavía."),
        ("2", "Normalización", "Unifica manual funcional y modelo en árbol en una única "
         "estructura, indexada por código CVN."),
        ("3", "Resolución auxiliar", "Resuelve las referencias del manual hacia entidades, "
         "tablas y tesauro, con evidencia registrada."),
        ("4", "Política semántica", "Decide, campo a campo, su forma semántica: texto, fecha, "
         "enum cerrado, catálogo abierto, referencia…"),
    ]
    n = len(steps)
    gap = 0.16
    cw = (11.9 - gap * (n - 1)) / n
    x = 0.7
    y = 2.05
    ch = 3.35
    for num, name, desc in steps:
        add_rect(s, x, y, cw, ch, fill=WHITE, line_color=LINE_SOFT, line_w=1.1, rounded=True, radius=0.08)
        circle_num(s, x + 0.22, y + 0.24, 0.55, num, fill=SECONDARY)
        add_text(s, x + 0.20, y + 1.0, cw - 0.4, 0.75, name, size=13.5, bold=True, color=DARK,
                  line_spacing=1.02)
        add_text(s, x + 0.20, y + 1.7, cw - 0.4, ch - 1.85, desc, size=10.8, color=TEXT_MUTED,
                  line_spacing=1.15)
        x += cw + gap

    add_text(s, 0.7, 5.65, 11.9, 0.5,
              "Etapas mantenidas a mano en src/cvn_codegen/, separadas de los bindings "
              "generados automáticamente en src/generated/.",
              size=11.5, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_notes(s,
              "Duracion objetivo: 100 s. Explicar las cuatro etapas del segundo bloque de la "
              "arquitectura: generacion estructural, normalizacion por codigo CVN, resolucion "
              "de referencias auxiliares y politica semantica. Insistir en que el codigo "
              "generado nunca se edita a mano: la limpieza semantica ocurre en capas "
              "posteriores.")


# ---------------------------------------------------------------------------
# Slide 11 - Resolucion de referencias auxiliares (cifras)
# ---------------------------------------------------------------------------

def slide_11_cifras_normalizacion():
    s = new_slide()
    kicker(s, "5 · Implementación del pipeline")
    title(s, "Normalización y resolución auxiliar, con evidencia")

    stats = [
        ("1457", "entradas normalizadas totales"),
        ("1429", "presentes en manual y en el modelo en árbol"),
        ("557", "declaran tabla de referencia en el manual"),
        ("33", "discrepancias manual / árbol registradas"),
        ("1", "referencia sin resolver (CVN_AGENCY_C)"),
    ]
    n = len(stats)
    gap = 0.18
    cw = (11.9 - gap * (n - 1)) / n
    x = 0.7
    y = 2.1
    for num, label in stats:
        stat_card(s, x, y, cw, 2.15, num, label,
                  num_color=ACCENT if num == "1" else PRIMARY)
        x += cw + gap

    add_rect(s, 0.7, 4.55, 11.9, 1.9, fill=CARD_BG, rounded=True, radius=0.09)
    add_bullets(s, 0.95, 4.75, 11.4, 1.55, [
        ("Evidencia, no suposición. ", "De las 557 entradas con tabla declarada, la resolución "
         "auxiliar se apoya en evidencia concreta del paquete; solo una queda sin resolver."),
        ("El caso CVN_AGENCY_C. ", "El manual cita una tabla que no tiene equivalente "
         "localizable en el resto del paquete: se conserva como referencia no resuelta, en "
         "lugar de forzar una interpretación."),
    ], size=12.5, space_after=8, line_spacing=1.05)

    add_notes(s,
              "Duracion objetivo: 80 s. Dar las cifras ya verificadas: 1457 entradas "
              "normalizadas, 1429 en ambas fuentes, 557 con tabla de referencia declarada, 33 "
              "discrepancias, 1 referencia no resuelta (CVN_AGENCY_C). El mensaje es que cada "
              "resolucion se apoya en evidencia concreta del paquete, no en suposiciones.")


# ---------------------------------------------------------------------------
# Slide 12 - Politica semantica
# ---------------------------------------------------------------------------

def slide_12_politica_semantica():
    s = new_slide()
    kicker(s, "5 · Implementación del pipeline")
    title(s, "Política semántica: nueve formas, con trazabilidad")

    forms = ["Texto", "Fecha", "Número", "Enum. cerrada", "Catálogo abierto",
              "Referencia a entidad", "Término de tesauro", "Subtipo", "Referencia no resuelta"]
    cols, rows_n = 3, 3
    gap = 0.14
    cw = (11.9 - gap * (cols - 1)) / cols
    ch = 0.62
    x0, y0 = 0.7, 1.95
    for i, f in enumerate(forms):
        r, c = divmod(i, cols)
        x = x0 + c * (cw + gap)
        y = y0 + r * (ch + gap)
        add_rect(s, x, y, cw, ch, fill=CARD_BG2, rounded=True, radius=0.20)
        add_text(s, x + 0.1, y, cw - 0.2, ch, f, size=13, bold=True, color=PRIMARY,
                  align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    code_y = y0 + rows_n * (ch + gap) + 0.10
    add_text(s, 0.7, code_y, 11.9, 0.35,
              "Ejemplo: campo de sexo (código 000.010.000.030, tabla CVN_SEX_A)",
              size=12.5, bold=True, color=DARK)
    add_rect(s, 0.7, code_y + 0.42, 11.9, 1.3, fill=DARK, rounded=True, radius=0.08)
    code_txt = (
        '{\n'
        '  "code": "000", "label": "Mujer",\n'
        '  "source": "CVN_SEX_A"\n'
        '}'
    )
    add_text(s, 1.0, code_y + 0.55, 6.5, 1.05, code_txt, size=13, color=TEXT_ON_DARK,
              font="Courier New", line_spacing=1.15)
    add_text(s, 7.7, code_y + 0.55, 4.6, 1.05,
              "Enumeración cerrada con solo dos valores posibles: forma semántica más simple, "
              "usada como ejemplo guía en toda la memoria.",
              size=11, color=TEXT_ON_DARK_MUTED, line_spacing=1.15)

    add_notes(s,
              "Duracion objetivo: 90 s. Enumerar rapido las nueve formas semanticas "
              "reconocidas. Detenerse en el ejemplo del campo de sexo (CVN_SEX_A): enumeracion "
              "cerrada con dos valores, el ejemplo mas simple, usado como hilo conductor en la "
              "memoria desde normalizacion hasta el JSON final.")


# ---------------------------------------------------------------------------
# Slide 13 - Modelos de dominio y modelo conceptual
# ---------------------------------------------------------------------------

def slide_13_modelos():
    s = new_slide()
    kicker(s, "5 · Implementación del pipeline")
    title(s, "Del modelo de dominio al modelo conceptual agnóstico")

    stats = [
        ("105", "archivos de modelos de dominio generados"),
        ("182", "definiciones en el JSON Schema resultante"),
        ("74", "de esas definiciones son vocabularios controlados"),
    ]
    x, y0, w, h, gap = 0.7, 2.0, 3.75, 1.32, 0.16
    y = y0
    for num, label in stats:
        stat_card(s, x, y, w, h, num, label, num_color=PRIMARY)
        y += h + gap

    add_text(s, x, y + 0.02, w, 0.85,
              "105 = 101 módulos por concepto curricular + 1 módulo solo-manual + 1 módulo de "
              "nodos de árbol + 1 módulo de 13 enumeraciones.",
              size=10.5, italic=True, color=TEXT_MUTED, line_spacing=1.15)

    add_image_fit(s, os.path.join(FIGS, "open_cvn_presentation_overview.png"),
                   4.75, 2.0, 7.85, 4.65)

    add_notes(s,
              "Duracion objetivo: 90 s. Explicar que esta etapa produce dos artefactos "
              "distintos: los modelos de dominio en Pydantic (105 archivos) y, por separado, "
              "el modelo conceptual agnostico que ya no vuelve a leer XML/XSD. Mostrar la "
              "vista UML compacta como el puente hacia el JSON Schema (182 definiciones, 74 "
              "vocabularios controlados).")


# ---------------------------------------------------------------------------
# Slide 14 - Formato Open CVN JSON
# ---------------------------------------------------------------------------

def slide_14_formato_json():
    s = new_slide()
    kicker(s, "6 · Formato Open CVN y herramienta")
    title(s, "El formato Open CVN JSON")

    add_rect(s, 0.7, 1.95, 6.0, 4.55, fill=DARK, rounded=True, radius=0.06)
    code_txt = (
        '{\n'
        '  "schema_version": "0.1.0",\n'
        '  "metadata": { "language": "es" },\n'
        '  "curriculum": {\n'
        '    "identity": {},\n'
        '    "education": [],\n'
        '    "research": [],\n'
        '    "professional_experience": [],\n'
        '    "achievements": [],\n'
        '    "other": []\n'
        '  },\n'
        '  "extensions": {}\n'
        '}'
    )
    add_text(s, 1.0, 2.20, 5.4, 4.1, code_txt, size=14, color=TEXT_ON_DARK,
              font="Courier New", line_spacing=1.28)

    add_bullets(s, 7.0, 1.95, 5.6, 4.6, [
        ("schema_version. ", "identifica la versión del formato, no la del JSON Schema."),
        ("metadata. ", "describe el documento, no su contenido curricular."),
        ("curriculum. ", "la información propiamente dicha, organizada por área conceptual del "
         "dominio académico."),
        ("extensions. ", "opcional, reservado para metadatos de una herramienta concreta."),
    ], size=14.5, space_after=14, line_spacing=1.1)

    add_text(s, 7.0, 5.85, 5.6, 0.75,
              "Los valores controlados comparten una misma forma: code, label, source y, "
              "opcionalmente, raw_value o uri.",
              size=11.5, italic=True, color=TEXT_MUTED, line_spacing=1.15)

    add_notes(s,
              "Duracion objetivo: 90 s. Mostrar la estructura raiz minima: cuatro campos "
              "(schema_version, metadata, curriculum, extensions). Explicar que curriculum se "
              "organiza por area conceptual (identity, education, research, professional "
              "experience, achievements, other) y no por estructura XML. Mencionar la forma "
              "comun de las referencias controladas, reutilizando el ejemplo de sexo ya visto.")


# ---------------------------------------------------------------------------
# Slide 15 - Validacion: parser y validador
# ---------------------------------------------------------------------------

def slide_15_validacion():
    s = new_slide()
    kicker(s, "6 · Formato Open CVN y herramienta")
    title(s, "Parser y validador: un contrato con cinco estados")

    states = [("Sin ejecutar", CARD_BG2, TEXT_DARK), ("Válido", SECONDARY, WHITE),
               ("Válido con avisos", ACCENT, WHITE), ("Inválido", RGBColor(0xB0, 0x3A, 0x2E), WHITE),
               ("Fallido", DARK, WHITE)]
    x, y0, w, h, gap = 0.7, 2.0, 5.55, 0.62, 0.14
    y = y0
    for label, fill, tc in states:
        add_rect(s, x, y, w, h, fill=fill, rounded=True, radius=0.5)
        add_text(s, x + 0.2, y, w - 0.4, h, label, size=13.5, bold=True, color=tc,
                  anchor=MSO_ANCHOR.MIDDLE)
        y += h + gap

    add_text(s, x, y + 0.06, w, 1.0,
              "Un aviso nunca convierte un documento válido en inválido: señala un caso "
              "semánticamente atípico (tipo inconsistente con la sección, referencia sin code "
              "ni label) para que pueda revisarse.",
              size=11.5, italic=True, color=TEXT_MUTED, line_spacing=1.15)

    add_image_fit(s, os.path.join(FIGS, "open_cvn_import_validation_flow.png"),
                   6.55, 1.9, 6.05, 4.75)

    add_notes(s,
              "Duracion objetivo: 90 s. El contrato publico expone cuatro funciones (leer, "
              "validar, importar XML, importar PDF) que devuelven siempre la misma estructura "
              "de resultado, con cinco estados posibles. Explicar por que existe 'valido con "
              "avisos': para senializar casos semanticamente atipicos sin bloquear el "
              "documento. Recorrer brevemente el diagrama de flujo de importacion/validacion.")


# ---------------------------------------------------------------------------
# Slide 16 - Importacion y herramienta local
# ---------------------------------------------------------------------------

def slide_16_importacion_herramienta():
    s = new_slide()
    kicker(s, "6 · Formato Open CVN y herramienta")
    title(s, "Importación determinista y herramienta local")

    paths = [
        ("1", "Open CVN JSON", "Caso directo: se analiza y valida sin transformación adicional."),
        ("2", "CVN XML", "Mapeo semántico parcial y creciente; lo no reconocido se preserva "
         "como diagnóstico."),
        ("3", "CVN PDF", "Prioriza el XML embebido; solo si falla, y con autorización "
         "explícita, recurre a un LLM opcional y trazable."),
    ]
    x, y0, w = 0.7, 1.95, 5.55
    y = y0
    for num, name, desc in paths:
        h = 1.28
        add_rect(s, x, y, w, h, fill=WHITE, line_color=LINE_SOFT, line_w=1.1, rounded=True, radius=0.10)
        circle_num(s, x + 0.18, y + 0.18, 0.5, num, fill=PRIMARY)
        add_text(s, x + 0.86, y + 0.14, w - 1.0, 0.35, name, size=14, bold=True, color=DARK)
        add_text(s, x + 0.86, y + 0.50, w - 1.0, h - 0.55, desc, size=10.8, color=TEXT_MUTED,
                  line_spacing=1.1)
        y += h + 0.16

    rows = [
        ["Comando", "Función"],
        ["store init", "Inicializa el almacenamiento SQLite local"],
        ["json import / pdf import", "Importa Open CVN JSON o un CVN PDF (LLM opcional)"],
        ["versions derive", "Crea una versión derivada a partir del maestro"],
        ["versions include / exclude", "Selecciona o excluye secciones y entradas"],
        ["latex export / pdf generate", "Exporta a LaTeX y compila a PDF"],
    ]
    tbl_shape = s.shapes.add_table(len(rows), 2, Inches(6.55), Inches(1.95), Inches(6.05), Inches(4.75))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.55)
    table.columns[1].width = Inches(3.50)
    fill_table(table, rows)
    style_table(table, len(rows), 2, body_font_size=11.5, header_font_size=12.5)

    add_notes(s,
              "Duracion objetivo: 100 s. Explicar los tres caminos de importacion en orden de "
              "prioridad: JSON directo, XML con mapeo parcial, PDF con extraccion determinista "
              "y LLM opcional solo como ultimo recurso autorizado y validado despues. Despues "
              "presentar la herramienta local: modelo maestro / version derivada y los "
              "comandos principales de la CLI.")


# ---------------------------------------------------------------------------
# Slide 17 - Evaluacion: 488 pruebas
# ---------------------------------------------------------------------------

def slide_17_evaluacion():
    s = new_slide()
    kicker(s, "7 · Evaluación y resultados")
    title(s, "488 pruebas automatizadas, sin ningún fallo")

    stats = [("488", "pruebas ejecutadas"), ("8", "niveles de evaluación"),
              ("692,80 s", "de ejecución paralela"), ("0", "fallos en la última ejecución")]
    n = len(stats)
    gap = 0.18
    cw = (11.9 - gap * (n - 1)) / n
    x = 0.7
    for num, label in stats:
        stat_card(s, x, 1.85, cw, 1.35, num, label, num_color=PRIMARY)
        x += cw + gap

    chart_data = CategoryChartData()
    chart_data.categories = ["Artefactos generados", "Normalización y trazabilidad",
                              "Política semántica", "JSON Schema y validación",
                              "Parsers e importadores", "Almacenamiento y versiones",
                              "Exportación", "Extremo a extremo"]
    chart_data.add_series("Pruebas", (146, 90, 76, 25, 63, 25, 19, 44))
    gframe = s.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(0.7), Inches(3.45),
                                 Inches(8.4), Inches(3.55), chart_data)
    chart = gframe.chart
    chart.has_legend = False
    chart.has_title = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0'
    plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = Pt(10)
    plot.data_labels.font.color.rgb = TEXT_DARK
    plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = PRIMARY
    cat_axis = chart.category_axis
    cat_axis.tick_labels.font.size = Pt(10.5)
    cat_axis.tick_labels.font.color.rgb = TEXT_DARK
    cat_axis.format.line.color.rgb = LINE_SOFT
    val_axis = chart.value_axis
    val_axis.visible = False
    val_axis.has_major_gridlines = False
    val_axis.minimum_scale = 0
    val_axis.maximum_scale = 165

    add_rect(s, 9.3, 3.45, 3.3, 3.55, fill=CARD_BG, rounded=True, radius=0.09)
    add_text(s, 9.55, 3.65, 2.85, 0.55, "uv run pytest -n auto tests", size=11.5, bold=True,
              color=PRIMARY, font="Courier New")
    add_bullets(s, 9.55, 4.35, 2.85, 2.5, [
        "Organizada por las mismas capas que la arquitectura.",
        "Se ejecuta también en integración continua en cada contribución.",
        "El resultado no depende del entorno de quien lo ejecuta.",
    ], size=11.5, space_after=10, line_spacing=1.12)

    add_notes(s,
              "Duracion objetivo: 100 s. Presentar el resultado central de la evaluacion: 488 "
              "pruebas, 8 niveles alineados con la arquitectura, 0 fallos, 692,80 segundos en "
              "paralelo. Un unico comando ejecuta todo, tanto en local como en CI. Este es el "
              "dato que sostiene todas las afirmaciones tecnicas anteriores.")


# ---------------------------------------------------------------------------
# Slide 18 - Discusion: garantias vs limitaciones
# ---------------------------------------------------------------------------

def slide_18_discusion():
    s = new_slide()
    kicker(s, "7 · Evaluación y resultados")
    title(s, "Garantías fuertes donde el proceso es determinista")

    add_rect(s, 0.7, 1.95, 5.85, 4.6, fill=TEAL_BG, rounded=True, radius=0.10)
    add_text(s, 0.95, 2.15, 5.35, 0.4, "GARANTÍAS FUERTES Y DETERMINISTAS", size=12.5, bold=True,
              color=SECONDARY)
    add_bullets(s, 0.95, 2.62, 5.4, 3.75, [
        "Generación estructural desde los esquemas XSD oficiales.",
        "Normalización y resolución de referencias con evidencia trazable.",
        "Política semántica evaluada a partir de evidencia concreta.",
        "Validación estructural y en tiempo de ejecución del formato Open CVN JSON.",
    ], size=13.5, space_after=16, bullet_hex="1C7293", line_spacing=1.12)

    add_rect(s, 6.75, 1.95, 5.85, 4.6, fill=AMBER_BG, rounded=True, radius=0.10)
    add_text(s, 7.0, 2.15, 5.35, 0.4, "GARANTÍAS PARCIALES Y DOCUMENTADAS", size=12.5, bold=True,
              color=ACCENT)
    add_bullets(s, 7.0, 2.62, 5.4, 3.75, [
        "Importación de CVN XML: mapeo semántico parcial y creciente.",
        "Importación asistida por LLM: opcional, no autoritativa sin revisión.",
        "Generación de PDF: depende de un motor LaTeX disponible en el entorno.",
    ], size=13.5, space_after=16, bullet_hex="E0993C", line_spacing=1.12)

    add_text(s, 0.7, 6.75, 11.9, 0.45,
              "No debe interpretarse como una conversión completa del ecosistema CVN ni como "
              "una validación semántica total del dominio.",
              size=11.5, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_notes(s,
              "Duracion objetivo: 90 s. Distinguir con precision: garantias fuertes donde el "
              "proceso es determinista (generacion, normalizacion, politica semantica, "
              "validacion Open CVN JSON) frente a garantias parciales y documentadas "
              "(importacion XML, importacion LLM, generacion PDF). Es el momento de dejar "
              "clara la honestidad tecnica del trabajo antes de pasar a conclusiones.")


# ---------------------------------------------------------------------------
# Slide 19 - Cumplimiento de objetivos
# ---------------------------------------------------------------------------

def slide_19_cumplimiento():
    s = new_slide()
    kicker(s, "8 · Conclusiones")
    title(s, "Cumplimiento de los objetivos específicos")

    rows = [
        ["Obj.", "Descripción resumida", "Estado"],
        ["OE1", "Análisis del ecosistema CVN y sus artefactos oficiales", "Cumplido"],
        ["OE2", "Estudio de alternativas de representación, modelado y validación", "Cumplido"],
        ["OE3", "Generación estructural reproducible desde XSD", "Cumplido"],
        ["OE4", "Normalización de metadatos CVN", "Cumplido"],
        ["OE5", "Definición de reglas semánticas trazables", "Cumplido"],
        ["OE6", "Generación de modelos de dominio y artefactos conceptuales", "Cumplido"],
        ["OE7", "Definición del formato Open CVN JSON", "Cumplido"],
        ["OE8", "Parser, validador, almacenamiento y exportación", "Cumplido"],
        ["OE9", "Importación determinista y asistida por LLM", "Parcial (deliberado)"],
        ["OE10", "Verificación automatizada y flujos reproducibles", "Cumplido"],
    ]
    tbl_shape = s.shapes.add_table(len(rows), 3, Inches(0.7), Inches(1.95), Inches(11.9), Inches(4.85))
    table = tbl_shape.table
    table.columns[0].width = Inches(0.95)
    table.columns[1].width = Inches(8.35)
    table.columns[2].width = Inches(2.6)
    fill_table(table, rows)
    style_table(table, len(rows), 3, body_font_size=12.5, header_font_size=13)
    for r in range(1, len(rows)):
        cell = table.cell(r, 2)
        parcial = "Parcial" in rows[r][2]
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.bold = True
                run.font.color.rgb = ACCENT if parcial else SECONDARY

    add_text(s, 0.7, 6.98, 11.9, 0.4,
              "El objetivo general se considera cumplido: la única salvedad es deliberada, no una carencia.",
              size=11.5, italic=True, color=TEXT_MUTED, align=PP_ALIGN.CENTER)

    add_notes(s,
              "Duracion objetivo: 90 s. Retomar los diez objetivos especificos del capitulo 1 "
              "en el mismo orden. Nueve se dan por cumplidos de forma determinista; OE9 "
              "(importacion) se cumple con garantias parciales de forma deliberada: la "
              "importacion asistida por LLM es opcional y no autoritativa sin revision, una "
              "decision de diseno, no una carencia.")


# ---------------------------------------------------------------------------
# Slide 20 - Contribuciones
# ---------------------------------------------------------------------------

def slide_20_contribuciones():
    s = new_slide()
    kicker(s, "8 · Conclusiones")
    title(s, "Siete contribuciones, demostradas capítulo a capítulo")

    items = [
        "Análisis computacional del ecosistema CVN",
        "Arquitectura reproducible por capas",
        "Trazabilidad entre CVN y Open CVN",
        "Formato Open CVN JSON",
        "Herramientas de validación y uso",
        "Importación determinista y asistida por LLM",
        "Verificación automatizada",
    ]
    cw = (11.9 - 3 * 0.16) / 4
    ch = 1.85
    row1 = items[:4]
    row2 = items[4:]
    y1 = 2.05
    x = 0.7
    for i, it in enumerate(row1):
        add_rect(s, x, y1, cw, ch, fill=CARD_BG, rounded=True, radius=0.10)
        circle_num(s, x + 0.18, y1 + 0.18, 0.48, i + 1, fill=PRIMARY)
        add_text(s, x + 0.18, y1 + 0.82, cw - 0.36, ch - 0.95, it, size=12.5, bold=True,
                  color=DARK, line_spacing=1.1)
        x += cw + 0.16
    y2 = y1 + ch + 0.2
    x2 = 0.7 + (11.9 - (3 * cw + 2 * 0.16)) / 2
    for i, it in enumerate(row2):
        add_rect(s, x2, y2, cw, ch, fill=CARD_BG, rounded=True, radius=0.10)
        circle_num(s, x2 + 0.18, y2 + 0.18, 0.48, i + 5, fill=SECONDARY)
        add_text(s, x2 + 0.18, y2 + 0.82, cw - 0.36, ch - 0.95, it, size=12.5, bold=True,
                  color=DARK, line_spacing=1.1)
        x2 += cw + 0.16

    add_notes(s,
              "Duracion objetivo: 70 s. Enumerar las siete contribuciones con ritmo, sin "
              "detenerse en cada una: analisis del ecosistema CVN, arquitectura reproducible "
              "por capas, trazabilidad CVN-Open CVN, formato JSON, herramientas de validacion, "
              "importacion determinista/LLM y verificacion automatizada. Recordar que ninguna "
              "sustituye al ecosistema CVN oficial: lo complementan.")


# ---------------------------------------------------------------------------
# Slide 21 - Conclusiones y trabajo futuro
# ---------------------------------------------------------------------------

def slide_21_conclusiones():
    s = new_slide()
    kicker(s, "8 · Conclusiones")
    title(s, "Conclusión y líneas de trabajo futuro")

    add_rect(s, 0.7, 1.95, 5.7, 4.6, fill=DARK, rounded=True, radius=0.09)
    add_text(s, 0.98, 2.2, 5.15, 0.4, "CONCLUSIÓN", size=12, bold=True, color=ACCENT)
    add_text(s, 0.98, 2.65, 5.15, 3.6,
              "Este trabajo demuestra que es posible construir, a partir de CVN, una "
              "arquitectura computacional abierta que separa la interoperabilidad "
              "estructural del significado curricular, con trazabilidad verificable y "
              "garantías reproducibles allí donde el proceso es determinista.\n\n"
              "Esa arquitectura por capas —más que el formato JSON o la herramienta que la "
              "demuestran— es la aportación principal de este Trabajo Fin de Grado.",
              size=14.5, italic=True, color=TEXT_ON_DARK, font=HEAD_FONT, line_spacing=1.22,
              space_after=6)

    future = [
        "Completar el mapeo semántico de importación XML con currículos reales",
        "Ampliar las validaciones de dominio fuera de la capa estructural generada",
        "Estudiar restricciones OCL sobre el modelo conceptual",
        "Reducir la dependencia de un motor LaTeX en la generación de PDF",
        "Explorar extensiones semánticas (JSON-LD) e integración institucional",
    ]
    x, y0, w = 6.75, 1.95, 5.85
    add_text(s, x, y0, w, 0.35, "TRABAJO FUTURO", size=12, bold=True, color=SECONDARY)
    y = y0 + 0.42
    for i, f in enumerate(future):
        h = 0.75
        add_rect(s, x, y, w, h, fill=CARD_BG, rounded=True, radius=0.14)
        circle_num(s, x + 0.15, y + 0.14, 0.46, i + 1, fill=SECONDARY, size=13)
        add_text(s, x + 0.72, y, w - 0.9, h, f, size=11.8, color=TEXT_DARK,
                  anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
        y += h + 0.13

    add_notes(s,
              "Duracion objetivo: 90 s. Cerrar con la conclusion central: la arquitectura por "
              "capas, no el JSON ni la herramienta, es la aportacion principal. Las "
              "limitaciones no comprometen esta conclusion: unas vienen del propio paquete "
              "CVN, otras son decisiones deliberadas de alcance de un TFG. Repasar rapido las "
              "cinco lineas futuras sin extenderse, dejando tiempo para preguntas.")


# ---------------------------------------------------------------------------
# Slide 22 - Cierre
# ---------------------------------------------------------------------------

def slide_22_cierre():
    s = new_slide(bg=DARK)

    layer_colors = [ACCENT, SECONDARY, PRIMARY]
    bar_y = 0.0
    widths = [5.6, 9.4, 13.333]
    for i, (w, col) in enumerate(zip(widths, layer_colors)):
        add_rect(s, 0, bar_y + i * 0.24, w, 0.20, fill=col)

    add_text(s, 0.7, 2.75, 11.9, 1.1, "Gracias por su atención", size=44, bold=True,
              color=WHITE, font=HEAD_FONT, align=PP_ALIGN.CENTER)
    add_text(s, 0.7, 3.85, 11.9, 0.6, "Quedo a su disposición para las preguntas del tribunal",
              size=18, italic=True, color=TEXT_ON_DARK, align=PP_ALIGN.CENTER)

    add_text(s, 0.7, 6.85, 11.9, 0.4,
              "Carlos Martínez Jaén · Open CVN · ESIIAB, Universidad de Castilla-La Mancha",
              size=11, color=TEXT_ON_DARK_MUTED, align=PP_ALIGN.CENTER)

    add_notes(s,
              "Cierre: agradecer al tribunal, indicar explicitamente que se abre el turno de "
              "preguntas segun la normativa (exponer objetivos, metodologia, contenido y "
              "conclusiones; despues, preguntas del tribunal). Mantener la calma, escuchar la "
              "pregunta completa antes de responder, y si hace falta pedir que se repita o "
              "aclare sin problema.")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    slide_01_portada()
    slide_02_indice()
    slide_03_contexto()
    slide_04_como_funciona()
    slide_05_objetivos()
    slide_06_limitaciones()
    slide_07_propuesta()
    slide_08_metodologia()
    slide_09_arquitectura()
    slide_10_pipeline()
    slide_11_cifras_normalizacion()
    slide_12_politica_semantica()
    slide_13_modelos()
    slide_14_formato_json()
    slide_15_validacion()
    slide_16_importacion_herramienta()
    slide_17_evaluacion()
    slide_18_discusion()
    slide_19_cumplimiento()
    slide_20_contribuciones()
    slide_21_conclusiones()
    slide_22_cierre()

    prs.save(OUT)
    print("Presentacion guardada en:", OUT)
    print("Total de diapositivas:", len(prs.slides))


if __name__ == "__main__":
    main()
